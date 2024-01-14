import pdftotext
import os
from nltk.tokenize import word_tokenize
from pathlib import Path
import json
import argparse
from pptx import Presentation
import nltk
nltk.download('punkt')

ROOT_DIRECTORY = os.path.dirname(os.path.realpath(__file__))

# Read PDF
def read_pdf(file_path):
    with open(file_path, 'rb') as file:
        #reader = PyPDF2.PdfReader(file)
        #text = "\n".join([page.extract_text() for page in reader.pages])
        pdf = pdftotext.PDF(file)
        text = "\n\n".join(pdf)
    print(text)
    return text

# Read PPTX
def read_pptx(file_path):
    with open(file_path, 'rb') as file:
        pptx = Presentation(file)
        shape_texts = []
        for slide in pptx.slides:
            # read the shapes in each slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    shape_texts.append(text)
        text = "\n\n".join(shape_texts)
    print(text)
    return text
    
def dummy_json(text):
    count = 1
    json_dict = {}
    for line in str.split(text, sep='\n'):
        words = word_tokenize(line)
        json = {'words': words, 'labels': ['O' for i in range(len(words))]}
        json_dict[count] = json
        count += 1
    return json_dict

if __name__ == '__main__':
    # CLI args
    parser = argparse.ArgumentParser()
    parser.add_argument('--pptx', default=None, action='store_true')
    parser.add_argument('--pdf', default=None, action='store_true')
    args = parser.parse_args()
    is_only_pptx_files, is_only_pdf_files = False, False
    if args.pptx != None:
        is_only_pptx_files = True
    if args.pdf != None:
        is_only_pdf_files = True
    
    directory = 'data'
    for dir in Path(directory).glob('*'):
        nwdir = os.path.join(ROOT_DIRECTORY, "json_slides", dir.stem)
        txdir = os.path.join(ROOT_DIRECTORY, "data_text", dir.stem)
        try:
            os.mkdir(nwdir)
        except:
            pass
        try:
            os.mkdir(txdir)
        except:
            pass

        if not is_only_pptx_files:
            pdf_files = Path(dir).glob('**/*.pdf')
            for file in pdf_files:
                print(file)
                text = read_pdf(file_path=file)
                with open(os.path.join(txdir, Path(file).stem)+".txt", 'w') as f:
                    f.write(text)
                dummy_json_text = dummy_json(text)
                with open(os.path.join(nwdir, Path(file).stem)+".json", 'w') as f:
                    f.write(json.dumps(dummy_json_text, indent=4))
        if not is_only_pdf_files:
            pptx_files = Path(dir).glob('**/*.pptx')
            for file in pptx_files:
                print(file)
                text = read_pptx(file_path=file)
                with open(os.path.join(txdir, Path(file).stem)+".txt", 'w') as f:
                    f.write(text)
                dummy_json_text = dummy_json(text)
                with open(os.path.join(nwdir, Path(file).stem)+".json", 'w') as f:
                    f.write(json.dumps(dummy_json_text, indent=4))
            


